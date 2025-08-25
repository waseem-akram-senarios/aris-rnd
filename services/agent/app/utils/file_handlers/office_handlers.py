"""Handlers for Microsoft Office file formats."""

import logging
from typing import List, Dict, Any
import io

from .base import BaseFileHandler, FileContent

logger = logging.getLogger(__name__)


class WordDocumentHandler(BaseFileHandler):
    """Handler for Word documents (.doc, .docx)."""
    
    SUPPORTED_EXTENSIONS = {'.doc', '.docx'}
    
    def can_handle(self, file_extension: str) -> bool:
        return file_extension.lower() in self.SUPPORTED_EXTENSIONS
    
    def extract_content(self, file_path: str, file_bytes: bytes) -> FileContent:
        """Extract content from Word document."""
        file_info = self.get_file_info(file_path)
        
        if not self.validate_file_size(file_bytes):
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error=f"File size exceeds {self.MAX_FILE_SIZE} bytes limit"
            )
        
        try:
            from docx import Document
            
            # Create Document from bytes
            doc = Document(io.BytesIO(file_bytes))
            
            text_parts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_text = [f"\n[Table {table_idx + 1}]"]
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        table_text.append(row_text)
                if len(table_text) > 1:  # Only add if table has content
                    text_parts.append("\n".join(table_text))
            
            text_content = "\n\n".join(text_parts)
            
            if not text_content.strip():
                return FileContent(
                    filename=file_info["filename"],
                    extension=file_info["extension"],
                    content_type="text",
                    text_content="[Document contains no extractable text]",
                    metadata={"paragraphs": 0, "tables": 0}
                )
            
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="text",
                text_content=text_content,
                metadata={
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(doc.tables)
                }
            )
            
        except Exception as e:
            self.logger.error(f"Error extracting Word document content: {str(e)}")
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error=str(e)
            )


class ExcelFileHandler(BaseFileHandler):
    """Handler for Excel files (.xls, .xlsx)."""
    
    SUPPORTED_EXTENSIONS = {'.xls', '.xlsx'}
    
    def can_handle(self, file_extension: str) -> bool:
        return file_extension.lower() in self.SUPPORTED_EXTENSIONS
    
    def extract_content(self, file_path: str, file_bytes: bytes) -> FileContent:
        """Extract content from Excel file."""
        file_info = self.get_file_info(file_path)
        
        if not self.validate_file_size(file_bytes):
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error=f"File size exceeds {self.MAX_FILE_SIZE} bytes limit"
            )
        
        try:
            import pandas as pd
            
            # Read Excel file from bytes
            excel_file = io.BytesIO(file_bytes)
            
            # Read all sheets
            all_sheets = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl' if file_info["extension"] == '.xlsx' else 'xlrd')
            
            text_parts = []
            sheet_count = len(all_sheets)
            total_rows = 0
            
            for sheet_name, df in all_sheets.items():
                text_parts.append(f"=== Sheet: {sheet_name} ===")
                
                if df.empty:
                    text_parts.append("[Empty sheet]")
                    continue
                
                # Limit rows to prevent overwhelming context
                max_rows = 50
                if len(df) > max_rows:
                    display_df = df.head(max_rows)
                    text_parts.append(display_df.to_string(index=False))
                    text_parts.append(f"\n... ({len(df) - max_rows} more rows)")
                else:
                    text_parts.append(df.to_string(index=False))
                
                total_rows += len(df)
                text_parts.append("")  # Add blank line between sheets
            
            text_content = "\n".join(text_parts)
            
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="table",
                text_content=text_content,
                metadata={
                    "sheets": sheet_count,
                    "total_rows": total_rows,
                    "sheet_names": list(all_sheets.keys())
                }
            )
            
        except ImportError as e:
            self.logger.error(f"Missing required library for Excel processing: {str(e)}")
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error="Excel processing libraries (pandas, openpyxl/xlrd) not installed"
            )
        except Exception as e:
            self.logger.error(f"Error extracting Excel content: {str(e)}")
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error=str(e)
            )


class PowerPointHandler(BaseFileHandler):
    """Handler for PowerPoint files (.ppt, .pptx)."""
    
    SUPPORTED_EXTENSIONS = {'.ppt', '.pptx'}
    
    def can_handle(self, file_extension: str) -> bool:
        return file_extension.lower() in self.SUPPORTED_EXTENSIONS
    
    def extract_content(self, file_path: str, file_bytes: bytes) -> FileContent:
        """Extract content from PowerPoint presentation."""
        file_info = self.get_file_info(file_path)
        
        if not self.validate_file_size(file_bytes):
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error=f"File size exceeds {self.MAX_FILE_SIZE} bytes limit"
            )
        
        try:
            from pptx import Presentation
            
            # Create Presentation from bytes
            prs = Presentation(io.BytesIO(file_bytes))
            
            text_parts = []
            slide_count = len(prs.slides)
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {slide_idx + 1} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        table_text = ["[Table]"]
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            if row_text.strip():
                                table_text.append(row_text)
                        if len(table_text) > 1:
                            slide_text.extend(table_text)
                
                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                    slide_text.append(f"[Notes: {slide.notes_slide.notes_text_frame.text.strip()}]")
                
                if len(slide_text) > 1:  # Only add if slide has content
                    text_parts.append("\n".join(slide_text))
            
            if not text_parts:
                return FileContent(
                    filename=file_info["filename"],
                    extension=file_info["extension"],
                    content_type="text",
                    text_content="[Presentation contains no extractable text]",
                    metadata={"slides": slide_count}
                )
            
            text_content = "\n\n".join(text_parts)
            
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="structured",
                text_content=text_content,
                metadata={"slides": slide_count}
            )
            
        except ImportError as e:
            self.logger.error(f"Missing required library for PowerPoint processing: {str(e)}")
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error="PowerPoint processing library (python-pptx) not installed"
            )
        except Exception as e:
            self.logger.error(f"Error extracting PowerPoint content: {str(e)}")
            return FileContent(
                filename=file_info["filename"],
                extension=file_info["extension"],
                content_type="error",
                text_content="",
                metadata={},
                error=str(e)
            )

